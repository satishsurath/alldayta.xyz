# This code takes all pdfs in documents, scans them, then generates paragraph embeddings on a sliding scale of 200 tokens (roughly 150 words)
# Run this before you run EmbedDocuments.py or app.py
# You need an OpenAI key saved in APIkey.txt
# Note that if your PDFs are not searchable, this won't work - use a third party tool to convert them to txt or doc first.  You
#   can look at the "-originaltext.csv" file created here and scan real quick to see if the text looks corrupted for any of your docs


import os
import time
from pdfminer.high_level import extract_text
import nltk
import pandas as pd
import numpy as np
import json
import io
# you need to pip install python-docx, not docx
import docx
from app import app
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


def read_settings(file_name):
    settings = {}
    with open(file_name, "r") as f:
        for line in f:
            key, value = line.strip().split("=")
            settings[key] = value
    return settings
#settings = read_settings("settings.txt")

def chunk_documents_given_course_name(course_name):
    try:
        nltk.data.find('tokenizers/punkt')
    except LookupError:
        nltk.download('punkt')


    # Set the desired chunk size and overlap size
    # chunk_size is how many tokens we will take in each block of text
    # overlap_size is how much overlap. So 200, 100 gives you chunks of between the 1st and 200th word, the 100th and 300th, the 200 and 400th...
    # I have in no way optimized these
    chunk_size = 200
    overlap_size = 100

    # load user settings and api key

    filedirectory = course_name
    # Check if the subfolder exists, if not, create it
    output_folder = os.path.join(course_name,"Textchunks")
    if not os.path.exists(output_folder):
        os.makedirs(output_folder)


    # Loop through all pdf, txt, tex in the "documents" folder
    for filename in os.listdir(filedirectory):
        if filename not in ["Textchunks.npy", "Textchunks-originaltext.csv", "CourseContentActivations.JSON", "course_meta.json"] and not filename.startswith('.'):
            filepath = os.path.join(filedirectory, filename)
            #Ingnore folders and hidden files
            if os.path.isfile(filepath) and not filename.startswith('.'):
                # Remove the file extension from the filename
                filename_without_extension = os.path.splitext(filename)[0]
                # Save the df_chunks to the output_folder subfolder with the new file name
                output_file = os.path.join(output_folder, filename_without_extension + "-originaltext.csv")
                # Create an empty DataFrame to store the text and title of each document
                df = pd.DataFrame(columns=["Title", "Text"])
                print("Loading " + filename)
                #check if filename already exits in the output folder with the filename as output_file and its not file.startswith("Syllabus" + course_name)
                if os.path.isfile(output_file) and not filename.startswith("Syllabus" + course_name):
                    print("File already exists. Skipping...")
                    app.logger.info(f"File already exists. Skipping... {filename}")
                    continue
                else:
                    print("File does not exist (or it maybe the Syllabus file.) Processing...")
                    app.logger.info(f"File does not exist. Processing... {filename}")
                    # 1. PDF
                    if filename.endswith(".pdf"):
                        app.logger.info(f"Processing PDF file: {filename}")
                        # Open the PDF file in read-binary mode
                        filepath = os.path.join(filedirectory, filename)
                        text  = extract_text(filepath)
                        # Add the text and title to the DataFrame
                        title = os.path.splitext(filename)[0]  # Remove the file extension from the filename
                        new_row = pd.DataFrame({"Title": [title], "Text": [text]})
                        df = pd.concat([df, new_row], ignore_index=True)


                    # 2. DOCX files
                    elif filename.endswith(".doc") or filename.endswith(".docx"):
                        app.logger.info(f"Processing DOCX file: {filename}")
                        # Open the DOC/DOCX file in binary mode and read the raw data
                        filepath = os.path.join(filedirectory, filename)
                        doc = docx.Document(filepath)

                        # Convert the file to UTF-8 and extract the text
                        text = ''
                        for paragraph in doc.paragraphs:
                            text += paragraph.text

                        # Add the text and title to the DataFrame
                        title = os.path.splitext(filename)[0]  # Remove the file extension from the filename
                        new_row = pd.DataFrame({"Title": [title], "Text": [text]})
                        df = pd.concat([df, new_row], ignore_index=True)



                    # 3. TXT files
                    elif filename.endswith(".txt"):
                        app.logger.info(f"Processing TXT file: {filename}")
                        # Open the text file and read its contents
                        filepath = os.path.join(filedirectory, filename)
                        with open(filepath, "r", encoding="utf-8") as file:
                            text = file.read()

                        # Add the text and title to the DataFrame
                        title = os.path.splitext(filename)[0]  # Remove the file extension from the filename
                        new_row = pd.DataFrame({"Title": [title], "Text": [text]})
                        df = pd.concat([df, new_row], ignore_index=True)
                        


                    # 4. LaTeX files
                    elif filename.endswith(".tex"):
                        app.logger.info(f"Processing LaTeX file: {filename}")
                        # Use regular expressions to extract regular text from the LaTeX file
                        filepath = os.path.join(filedirectory, filename)
                        with open(filepath, "r", encoding="utf-8") as file:
                            text = file.read()
                        
                        # Add the text and title to the DataFrame
                        title = os.path.splitext(filename)[0] # Remove the file extension from the filename
                        new_row = pd.DataFrame({"Title": [title], "Text": [text]})
                        df = pd.concat([df, new_row], ignore_index=True)

                        

                    elif filename.endswith(".pptx") or filename.endswith(".ppt"):
                        app.logger.info(f"Processing PPTX file: {filename}")
                        filepath = os.path.join(filedirectory, filename)
                        prs = Presentation(filepath)
                        # text_runs will be populated with a list of strings,
                        # one for each text run in presentation
                        text_runs = []
                        text = ''
                        for slide in prs.slides:
                            for shape in slide.shapes:
                                if not shape.has_text_frame:
                                    continue
                                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                                    continue                    
                                for paragraph in shape.text_frame.paragraphs:
                                    for run in paragraph.runs:
                                        run_text = run.text
                                        if run_text and run_text[-1] != '.':
                                            run_text += '.'  # Add a period at the end if it is not already there
                                        text_runs.append(run_text)
                                        text += run_text
                        title = os.path.splitext(filename)[0] # Remove the file extension from the filename
                        new_row = pd.DataFrame({"Title": [title], "Text": [text]})
                        df = pd.concat([df, new_row], ignore_index=True)

                        
                    ## Backend Idea 1:
                    # Instead of Tokenizing each row of text on its own
                    # Consider first splitting them into an "Array of sentences" by splitting them with the "." character
                    # This way the text rows will not be split mid-sentence
                    # Only split sentences that are "bigger" than chuck size      
                        
                    # Loop through the rows and create overlapping chunks for each text
                    chunks = []
                    for i, row in df.iterrows():
                        # Tokenize the text for the current row
                        tokens = nltk.word_tokenize(row['Text'])

                        # Loop through the tokens and create overlapping chunks
                        for j in range(0, len(tokens), chunk_size - overlap_size):
                            # Get the start and end indices of the current chunk
                            start = j
                            end = j + chunk_size

                            # Create the current chunk by joining the tokens within the start and end indices
                            chunk = ' '.join(tokens[start:end])

                            # Add the article title to the beginning of the chunk
                            chunk_with_title = "This text comes from the document " + row['Title'] + ". " + chunk

                            # Append the current chunk to the list of chunks, along with the corresponding title
                            chunks.append([row['Title'], chunk_with_title])

                    # Convert the list of chunks to a dataframe
                    df_chunks = pd.DataFrame(chunks, columns=['Title', 'Text'])

                    # Truncate the filename if it's too long, e.g., limit to 250 characters
                    max_filename_length = 250
                    if len(filename) > max_filename_length:
                        filename = filename[:max_filename_length]
                    
                    df_chunks.to_csv(output_file, encoding='utf-8', escapechar='\\', index=False)

                    print("Saving " + filename)
                    app.logger.info(f"Saved {filename} to '{output_file}'")





